from pptx import Presentation
from pptx.util import Inches
import Dra
import time

class Powerpoint:
    def __init__(self):
        pass

    def gantt_preso(self, data):
        # Data is a list of lists.

        # Create presentation and slide.
        prs = Presentation("~/Ref/Theme.pptx")
        title_slide_layout = prs.slide_layouts[1]

        for x in range(len(data)):
            slide = prs.slides.add_slide(title_slide_layout)

            # Create title
            title = slide.shapes.title
            title.text = str(data[x][0][3]) + " Campaign - " \
            + str(data[x][0][10]) + " - " + str(data[x][0][9])

            # Create image
            img = Dra.Gantt(data[x])
            img_path = "~/Final/" + str(img.filename)

            # Add image to slide
            left = Inches(0)
            top = Inches(1.25)
            width = Inches(10)
            pic = slide.shapes.add_picture(img_path, left, top, width=width)
	
        # Save it all out
        prs.save("~/Final/Campaign Charts_" + str(time.strftime("%b%d")) \
        + "_" + str(data[x][0][9]) + ".pptx")
